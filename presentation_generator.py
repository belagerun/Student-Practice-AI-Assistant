import re
from datetime import datetime
from pathlib import Path


ARTIFACTS_DIR = Path(__file__).resolve().with_name("artifacts")


def _clean_text(text):
    return " ".join((text or "").split())


def sanitize_presentation_outline(content):
    forbidden_patterns = [
        r"^\s*вот\s+презентац.*$",
        r"^\s*ниже\s+представлен.*$",
        r"^\s*ниже\s+приведен.*$",
        r"^\s*ниже\s+приведён.*$",
        r"^\s*here\s+is\s+.*presentation.*$",
        r"^\s*below\s+is\s+.*presentation.*$",
        r"^\s*this\s+is\s+.*presentation.*$",
        r"^\s*generated\s+by.*$",
        r"^\s*запрос\s+пользователя.*$",
        r"^\s*user\s+request.*$",
        r"^\s*source\s+material.*$",
        r"^\s*конечно[,.!]?.*$",
    ]
    cleaned_lines = []

    for raw_line in (content or "").splitlines():
        line = raw_line.strip()

        if not line:
            cleaned_lines.append("")
            continue

        lowered = line.lower()

        if any(re.match(pattern, lowered) for pattern in forbidden_patterns):
            continue

        if re.fullmatch(r"(?:slide|слайд)\s*\d+", lowered):
            continue

        line = re.sub(
            r"^\s*(?:slide|слайд)\s*(\d+)\s*[-.]\s*",
            r"SLIDE \1: ",
            line,
            flags=re.IGNORECASE,
        )
        line = re.sub(
            r"^\s*(?:slide|слайд)\s*(\d+)\s*:\s*(?:slide|слайд)\s*\1\s*$",
            r"SLIDE \1:",
            line,
            flags=re.IGNORECASE,
        )
        line = re.sub(
            r"^\s*(?:slide|слайд)\s*(\d+)\s*:\s*(?:slide|слайд)\s*\d+\s*[-:]?\s*",
            r"SLIDE \1: ",
            line,
            flags=re.IGNORECASE,
        )
        cleaned_lines.append(line)

    cleaned = "\n".join(cleaned_lines).strip()
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)

    return cleaned


def _split_into_sections(content, slides_count):
    content = sanitize_presentation_outline(content)

    if not content:
        raise ValueError("Presentation content is empty.")

    slide_blocks = re.split(
        r"(?:^|\n)\s*(?:SLIDE|Slide|СЛАЙД|Слайд)\s*\d*\s*[:.-]?\s*",
        content,
    )
    slide_blocks = [
        block.strip()
        for block in slide_blocks
        if block.strip()
    ]
    slide_blocks = [
        block
        for block in slide_blocks
        if block.lower() not in {"slide", "слайд"}
    ]

    if len(slide_blocks) >= 2:
        return slide_blocks[:slides_count]

    paragraphs = [
        paragraph.strip()
        for paragraph in re.split(r"\n{2,}", content)
        if paragraph.strip()
    ]

    if len(paragraphs) >= slides_count - 2:
        return paragraphs[:slides_count]

    sentences = re.split(r"(?<=[.!?])\s+", _clean_text(content))
    sentences = [
        sentence.strip()
        for sentence in sentences
        if sentence.strip()
    ]

    body_slides = max(1, slides_count - 2)
    chunk_size = max(1, len(sentences) // body_slides)
    sections = []

    for index in range(0, len(sentences), chunk_size):
        sections.append(" ".join(sentences[index:index + chunk_size]))

        if len(sections) >= body_slides:
            break

    return sections or [content]


def _extract_title_and_bullets(section, fallback_title):
    lines = [
        line.strip(" -*\t")
        for line in section.splitlines()
        if line.strip()
    ]

    if not lines:
        return fallback_title, []

    title = lines[0]
    body_lines = lines[1:]

    if re.fullmatch(r"(?:slide|слайд)\s*\d*", title, flags=re.IGNORECASE):
        title = fallback_title

    if not body_lines:
        body_lines = re.split(r"(?<=[.!?])\s+", lines[0])
        title = fallback_title

    bullets = []

    for line in body_lines:
        parts = re.split(r";\s+|(?<=[.!?])\s+", line)

        for part in parts:
            item = _clean_text(part.strip(" -*\t"))

            if item and item.lower() != title.lower():
                bullets.append(item[:180])

            if len(bullets) >= 7:
                break

        if len(bullets) >= 7:
            break

    return title[:80], bullets[:7]


def _add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    if not bullets:
        bullets = ["Key point"]

    for index, bullet in enumerate(bullets[:7]):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0


def generate_presentation(title, content, slides_count=10):
    try:
        from pptx import Presentation
    except ImportError as error:
        raise RuntimeError(
            "python-pptx is not installed. Add python-pptx to requirements.txt."
        ) from error

    title = _clean_text(title) or "Presentation"
    slides_count = max(3, min(int(slides_count or 10), 20))
    sections = _split_into_sections(content, slides_count)

    ARTIFACTS_DIR.mkdir(
        parents=True,
        exist_ok=True
    )

    prs = Presentation()
    title_section = sections[0] if sections else title
    title_slide_title, title_slide_bullets = _extract_title_and_bullets(
        title_section,
        title
    )
    subtitle = title_slide_bullets[0] if title_slide_bullets else "Краткий обзор"

    _add_title_slide(
        prs,
        title_slide_title,
        subtitle
    )

    body_sections = sections[1:slides_count - 1]

    for index, section in enumerate(body_sections, start=1):
        slide_title, bullets = _extract_title_and_bullets(
            section,
            f"Section {index}"
        )

        if not slide_title or slide_title.lower().startswith("slide "):
            slide_title = f"Section {index}"

        _add_bullet_slide(
            prs,
            slide_title,
            bullets
        )

    _add_bullet_slide(
        prs,
        "Conclusion / Thank You",
        [
            "Main ideas summarized",
            "Questions and discussion",
            "Thank you for your attention",
        ]
    )

    file_name = (
        "presentation_"
        f"{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    )
    file_path = ARTIFACTS_DIR / file_name
    prs.save(file_path)

    return {
        "file_path": str(file_path),
        "slides_count": len(prs.slides),
        "file_size": file_path.stat().st_size,
    }
